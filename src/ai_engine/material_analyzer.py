"""
Material Analyzer - Advanced analysis of lecture notes, slides, and external resources
Extracts key insights, concepts, and themes from uploaded materials
"""

import os
import re
import json
from typing import Dict, List, Tuple, Optional, Any
from pathlib import Path
from collections import Counter
import logging

logger = logging.getLogger(__name__)


class MaterialAnalyzer:
    """
    Analyze lecture materials (PDFs, slides, notes, resources) to extract:
    - Key concepts and topics
    - Learning objectives
    - Important definitions
    - Structure and hierarchy
    - Main themes and connections
    - Difficulty level
    - Recommended focus areas
    """

    def __init__(self):
        """Initialize material analyzer."""
        self.max_keywords = 20
        self.min_keyword_length = 3
        self.concept_markers = [
            "define", "definition", "is", "are", "means", "refers to",
            "concept", "term", "principle", "law", "theory", "model"
        ]
        self.objective_markers = [
            "learn", "understand", "analyze", "evaluate", "create",
            "students will", "you will", "upon completion", "objective",
            "goal", "learning outcome"
        ]

    def analyze_material(self, content: str, filename: str = "") -> Dict[str, Any]:
        """
        Comprehensive analysis of uploaded material.

        Args:
            content: Text content extracted from material
            filename: Original filename for context

        Returns:
            Dictionary containing:
                - key_concepts: Important topics and concepts
                - learning_objectives: Main learning goals
                - key_definitions: Important definitions found
                - structure: Document structure analysis
                - main_themes: Primary themes and topics
                - difficulty_level: Estimated difficulty (beginner/intermediate/advanced)
                - content_type: Type of material (lecture, slides, notes, etc.)
                - summary: Brief overview
                - focus_areas: Recommended areas to focus on
                - metadata: Document metadata and statistics
        """
        if not content or len(content.strip()) < 50:
            return self._empty_analysis()

        try:
            analysis = {
                "key_concepts": self._extract_concepts(content),
                "learning_objectives": self._extract_objectives(content),
                "key_definitions": self._extract_definitions(content),
                "structure": self._analyze_structure(content),
                "main_themes": self._extract_themes(content),
                "difficulty_level": self._estimate_difficulty(content),
                "content_type": self._identify_content_type(content, filename),
                "summary": self._generate_summary(content),
                "focus_areas": self._identify_focus_areas(content),
                "metadata": self._extract_metadata(content, filename),
            }

            logger.info(f"Material analysis complete. Concepts found: {len(analysis['key_concepts'])}")
            return analysis

        except Exception as e:
            logger.error(f"Error analyzing material: {str(e)}")
            return self._empty_analysis()

    def _extract_concepts(self, content: str) -> List[Dict[str, Any]]:
        """
        Extract key concepts and topics from content.

        Returns:
            List of concepts with importance scores
        """
        # Remove URLs and special characters
        clean_content = re.sub(r'http\S+|[^\w\s]', ' ', content.lower())
        words = clean_content.split()

        # Filter by length and common words
        stop_words = {
            'the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for',
            'of', 'with', 'is', 'are', 'be', 'was', 'were', 'that', 'this',
            'from', 'by', 'as', 'it', 'you', 'he', 'she', 'we', 'they'
        }

        filtered_words = [
            w for w in words
            if len(w) >= self.min_keyword_length and w not in stop_words
        ]

        # Count frequencies
        word_freq = Counter(filtered_words)

        # Identify noun phrases (technical terms)
        technical_terms = self._extract_technical_terms(content)

        # Combine and rank
        concepts = []
        all_concepts = dict(word_freq.most_common(self.max_keywords))

        # Add technical terms with higher weight
        for term in technical_terms[:self.max_keywords]:
            term_lower = term.lower()
            if term_lower not in all_concepts:
                all_concepts[term_lower] = len(technical_terms) / (technical_terms.index(term) + 1)

        # Create concept list with scores
        for concept, frequency in sorted(all_concepts.items(), key=lambda x: x[1], reverse=True)[:self.max_keywords]:
            concepts.append({
                "concept": concept.title(),
                "frequency": frequency,
                "importance": min(100, int((frequency / max(all_concepts.values())) * 100))
            })

        return concepts

    def _extract_technical_terms(self, content: str) -> List[str]:
        """
        Extract capitalized technical terms and proper nouns.
        """
        # Find words that are capitalized (likely technical terms or proper nouns)
        pattern = r'\b[A-Z][a-z]+(?:\s+[A-Z][a-z]+)*\b'
        terms = re.findall(pattern, content)

        # Count and return most common
        term_freq = Counter(terms)
        return [term for term, _ in term_freq.most_common(20)]

    def _extract_objectives(self, content: str) -> List[str]:
        """
        Extract learning objectives from content.
        """
        objectives = []
        sentences = re.split(r'[.!?]+', content)

        for sentence in sentences:
            sentence_lower = sentence.lower().strip()

            # Check for objective markers
            for marker in self.objective_markers:
                if marker in sentence_lower:
                    # Extract meaningful objectives
                    obj = self._clean_objective(sentence.strip())
                    if len(obj) > 20 and len(obj) < 200:
                        objectives.append(obj)
                    break

        # Remove duplicates while preserving order
        seen = set()
        unique_objectives = []
        for obj in objectives:
            if obj not in seen and len(unique_objectives) < 10:
                seen.add(obj)
                unique_objectives.append(obj)

        return unique_objectives

    def _clean_objective(self, text: str) -> str:
        """Clean and format objective text."""
        # Remove extra whitespace
        text = re.sub(r'\s+', ' ', text).strip()
        # Remove common markers
        text = re.sub(r'^(to|the|by|after|upon completion)?\s*', '', text, flags=re.IGNORECASE)
        return text

    def _extract_definitions(self, content: str) -> List[Dict[str, str]]:
        """
        Extract key definitions from content.
        """
        definitions = []
        sentences = re.split(r'[.!?]+', content)

        for sentence in sentences:
            sentence_lower = sentence.lower().strip()

            # Look for definition patterns
            for marker in self.concept_markers:
                if marker in sentence_lower and len(sentence) > 30:
                    # Extract term and definition
                    parts = re.split(rf'\b{marker}\b', sentence, maxsplit=1, flags=re.IGNORECASE)
                    if len(parts) == 2:
                        term = self._extract_term(parts[0])
                        definition = parts[1].strip()

                        if term and len(definition) > 20 and len(definition) < 300:
                            definitions.append({
                                "term": term,
                                "definition": definition
                            })
                    break

        # Keep unique definitions (up to 15)
        seen = set()
        unique_defs = []
        for d in definitions:
            if d['term'] not in seen and len(unique_defs) < 15:
                seen.add(d['term'])
                unique_defs.append(d)

        return unique_defs

    def _extract_term(self, text: str) -> Optional[str]:
        """Extract the term from definition context."""
        # Look for quoted text or emphasized text
        match = re.search(r'["\']([^"\']+)["\']|(?:^|\s)(\b[A-Z][a-z]+(?:\s+[A-Z][a-z]+)*\b)(?:\s|$)', text)
        if match:
            return match.group(1) or match.group(2)
        return None

    def _analyze_structure(self, content: str) -> Dict[str, Any]:
        """
        Analyze document structure.
        """
        lines = content.split('\n')
        headers = [line for line in lines if line.startswith('#') or (len(line) < 100 and line.isupper())]
        paragraphs = [line for line in lines if len(line.strip()) > 50]

        return {
            "total_lines": len(lines),
            "total_paragraphs": len(paragraphs),
            "estimated_sections": len(headers),
            "average_paragraph_length": int(sum(len(p) for p in paragraphs) / len(paragraphs)) if paragraphs else 0,
            "has_lists": bool(re.search(r'^\s*[-•*]\s', content, re.MULTILINE)),
            "has_numbering": bool(re.search(r'^\s*\d+[.)]\s', content, re.MULTILINE)),
        }

    def _extract_themes(self, content: str) -> List[Dict[str, Any]]:
        """
        Extract main themes and topics.
        """
        # Use concepts as themes but with different filtering
        concepts = self._extract_concepts(content)

        # Group related concepts into themes
        themes = []
        for concept in concepts[:10]:  # Top 10 concepts become themes
            related_count = len(re.findall(rf'\b{concept["concept"]}\b', content, re.IGNORECASE))

            themes.append({
                "theme": concept["concept"],
                "mentions": related_count,
                "importance": concept["importance"]
            })

        return themes

    def _estimate_difficulty(self, content: str) -> str:
        """
        Estimate difficulty level based on content analysis.
        """
        # Simple heuristics
        avg_word_length = sum(len(w) for w in content.split()) / max(len(content.split()), 1)
        technical_terms = len(self._extract_technical_terms(content))
        complex_words = len(re.findall(r'\b\w{10,}\b', content))

        score = (avg_word_length - 4) + (technical_terms / 5) + (complex_words / 50)

        if score > 15:
            return "Advanced"
        elif score > 10:
            return "Intermediate"
        else:
            return "Beginner"

    def _identify_content_type(self, content: str, filename: str = "") -> str:
        """
        Identify type of material.
        """
        content_lower = content.lower()

        # Check filename
        if filename:
            filename_lower = filename.lower()
            if 'slide' in filename_lower or 'presentation' in filename_lower:
                return "Presentation Slides"
            elif 'note' in filename_lower or 'lecture' in filename_lower:
                return "Lecture Notes"
            elif 'assignment' in filename_lower or 'exercise' in filename_lower:
                return "Assignment/Exercise"

        # Check content markers
        if re.search(r'slide\s*\d+|^\s*slide:', content, re.IGNORECASE):
            return "Presentation Slides"
        elif re.search(r'objective|learning outcomes|upon completion', content, re.IGNORECASE):
            return "Lecture Notes"
        elif re.search(r'question|problem|exercise|assignment', content, re.IGNORECASE):
            return "Assignment/Exercise"
        elif re.search(r'reference|bibliography|citation', content, re.IGNORECASE):
            return "Reference Material"
        else:
            return "General Material"

    def _generate_summary(self, content: str) -> str:
        """
        Generate brief summary of material.
        """
        # Extract first meaningful paragraph
        paragraphs = [p.strip() for p in content.split('\n') if len(p.strip()) > 50]

        if paragraphs:
            summary = paragraphs[0]
            # Limit to 150 characters
            if len(summary) > 150:
                summary = summary[:150].rsplit(' ', 1)[0] + "..."
            return summary
        return "No summary available"

    def _identify_focus_areas(self, content: str) -> List[str]:
        """
        Identify areas that students should focus on.
        """
        focus_areas = []

        # Check for emphasis markers
        emphasized = re.findall(r'\*\*([^*]+)\*\*|__([^_]+)__', content)
        for item in emphasized[:5]:
            term = item[0] or item[1]
            if len(term) > 5:
                focus_areas.append(f"Focus on: {term}")

        # Check for repeated concepts
        concepts = self._extract_concepts(content)
        for concept in concepts[:3]:
            if concept['frequency'] > 2:
                focus_areas.append(f"Important concept: {concept['concept']}")

        # Check for difficult sections
        if self._estimate_difficulty(content) == "Advanced":
            focus_areas.append("This material contains advanced topics - review fundamentals first")

        return focus_areas if focus_areas else ["Review all key concepts thoroughly"]

    def _extract_metadata(self, content: str, filename: str = "") -> Dict[str, Any]:
        """
        Extract metadata about the material.
        """
        words = content.split()
        sentences = re.split(r'[.!?]+', content)

        return {
            "total_words": len(words),
            "total_sentences": len(sentences),
            "avg_sentence_length": len(words) / max(len(sentences), 1),
            "unique_words": len(set(w.lower() for w in words)),
            "filename": filename or "Unknown",
            "content_length_category": self._categorize_length(len(content)),
        }

    def _categorize_length(self, length: int) -> str:
        """Categorize content by length."""
        if length < 1000:
            return "Short (< 1KB)"
        elif length < 5000:
            return "Medium (1-5KB)"
        elif length < 20000:
            return "Long (5-20KB)"
        else:
            return "Very Long (> 20KB)"

    def _empty_analysis(self) -> Dict[str, Any]:
        """Return empty analysis structure."""
        return {
            "key_concepts": [],
            "learning_objectives": [],
            "key_definitions": [],
            "structure": {
                "total_lines": 0,
                "total_paragraphs": 0,
                "estimated_sections": 0,
                "average_paragraph_length": 0,
                "has_lists": False,
                "has_numbering": False,
            },
            "main_themes": [],
            "difficulty_level": "Unknown",
            "content_type": "Unknown",
            "summary": "No content to analyze",
            "focus_areas": [],
            "metadata": {
                "total_words": 0,
                "total_sentences": 0,
                "avg_sentence_length": 0,
                "unique_words": 0,
                "filename": "",
                "content_length_category": "Empty",
            },
        }

    def compare_materials(self, analysis_list: List[Dict[str, Any]]) -> Dict[str, Any]:
        """
        Compare multiple materials to identify gaps, overlaps, and complementarity.

        Args:
            analysis_list: List of material analyses

        Returns:
            Comparison results
        """
        if not analysis_list:
            return {}

        # Extract all concepts
        all_concepts = []
        for analysis in analysis_list:
            all_concepts.extend([c['concept'] for c in analysis.get('key_concepts', [])])

        concept_freq = Counter(all_concepts)
        shared_concepts = [c for c, freq in concept_freq.items() if freq > 1]

        # Extract all objectives
        all_objectives = []
        for analysis in analysis_list:
            all_objectives.extend(analysis.get('learning_objectives', []))

        return {
            "shared_concepts": shared_concepts,
            "unique_concepts_per_material": [
                len(analysis.get('key_concepts', []))
                for analysis in analysis_list
            ],
            "total_unique_concepts": len(set(all_concepts)),
            "coverage_analysis": {
                "highly_covered": [c for c, freq in concept_freq.items() if freq == len(analysis_list)],
                "partially_covered": [c for c, freq in concept_freq.items() if 1 < freq < len(analysis_list)],
            },
            "total_objectives": len(all_objectives),
            "material_count": len(analysis_list),
        }


class MaterialProcessor:
    """
    Process and prepare materials for analysis and content generation.
    """

    def __init__(self):
        """Initialize material processor."""
        self.analyzer = MaterialAnalyzer()
        self.supported_formats = ['.pdf', '.docx', '.txt', '.md', '.doc', '.pptx']

    def process_material(self, file_path: str) -> Tuple[Dict[str, Any], str]:
        """
        Process uploaded material file.

        Args:
            file_path: Path to uploaded file

        Returns:
            Tuple of (analysis, extracted_content)
        """
        try:
            # Extract content based on file type
            content, filename = self._extract_content(file_path)

            # Analyze content
            analysis = self.analyzer.analyze_material(content, filename)

            return analysis, content

        except Exception as e:
            logger.error(f"Error processing material: {str(e)}")
            return self.analyzer._empty_analysis(), ""

    def _extract_content(self, file_path: str) -> Tuple[str, str]:
        """
        Extract content from various file formats.

        Returns:
            Tuple of (content, filename)
        """
        from pathlib import Path

        file_ext = Path(file_path).suffix.lower()
        filename = Path(file_path).name

        if file_ext == '.pdf':
            return self._extract_pdf(file_path), filename
        elif file_ext in ['.docx', '.doc']:
            return self._extract_word(file_path), filename
        elif file_ext == '.pptx':
            return self._extract_powerpoint(file_path), filename
        elif file_ext in ['.txt', '.md']:
            return self._extract_text(file_path), filename
        else:
            return self._extract_text(file_path), filename

    def _extract_pdf(self, file_path: str) -> str:
        """Extract text from PDF."""
        try:
            import pdfplumber

            content = []
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        content.append(text)

            return "\n\n".join(content)
        except ImportError:
            logger.warning("pdfplumber not installed, attempting fallback")
            return ""
        except Exception as e:
            logger.error(f"PDF extraction error: {str(e)}")
            return ""

    def _extract_word(self, file_path: str) -> str:
        """Extract text from Word document."""
        try:
            from docx import Document

            doc = Document(file_path)
            content = [para.text for para in doc.paragraphs if para.text.strip()]
            return "\n\n".join(content)
        except ImportError:
            logger.warning("python-docx not installed")
            return ""
        except Exception as e:
            logger.error(f"Word extraction error: {str(e)}")
            return ""

    def _extract_powerpoint(self, file_path: str) -> str:
        """Extract text from PowerPoint presentation."""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            content = []

            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                if slide_text:
                    content.append("SLIDE: " + " | ".join(slide_text))

            return "\n\n".join(content)
        except ImportError:
            logger.warning("python-pptx not installed")
            return ""
        except Exception as e:
            logger.error(f"PowerPoint extraction error: {str(e)}")
            return ""

    def _extract_text(self, file_path: str) -> str:
        """Extract text from plain text files."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            logger.error(f"Text extraction error: {str(e)}")
            return ""
